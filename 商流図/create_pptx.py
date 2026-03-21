# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# RGBColorの正しいインポート
from pptx.dml.color import RGBColor

# A4横サイズのプレゼンテーション作成
prs = Presentation()
prs.slide_width = Cm(29.7)
prs.slide_height = Cm(21.0)

# 空白レイアウトを使用
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# 色定義
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
ORANGE = RGBColor(255, 153, 0)
GREEN = RGBColor(0, 128, 0)
GRAY = RGBColor(128, 128, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=None, align=PP_ALIGN.CENTER):
    """テキストボックスを追加"""
    shape = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = color
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, text="", font_size=11, font_color=None):
    """角丸四角形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(100, 100, 100)

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if font_color:
            tf.paragraphs[0].font.color.rgb = font_color
        # 垂直中央揃え
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape

# ===== タイトル =====
title_shape = add_text_box(slide, 1, 0.5, 27.7, 1.5,
    "FIT発電所向けPCS交換ソリューション事業　商流図",
    font_size=22, bold=True, color=DARK_BLUE)

# ===== サブタイトル =====
add_text_box(slide, 1, 1.8, 27.7, 0.8,
    "～ ユアサ商事 × オルテナジー × ダイヤゼブラ電機 の3社協業スキーム ～",
    font_size=12, color=GRAY)

# ===== 背景説明ボックス =====
bg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(2.8), Cm(27.7), Cm(2.2))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = RGBColor(230, 242, 255)
bg_box.line.color.rgb = LIGHT_BLUE

bg_text = bg_box.text_frame
bg_text.word_wrap = True
p1 = bg_text.paragraphs[0]
p1.text = "【背景】"
p1.font.size = Pt(11)
p1.font.bold = True
p1.font.color.rgb = DARK_BLUE

p2 = bg_text.add_paragraph()
p2.text = "ユアサ商事がダイヤゼブラ社製PCSを納品してから約10年が経過。FIT発電所オーナーからPCS交換を含む修繕依頼が殺到。"
p2.font.size = Pt(10)

p3 = bg_text.add_paragraph()
p3.text = "オルテナジーが「パワまる」サービスを通じてダイヤゼブラ製最新PCSへの交換ソリューションを提供。"
p3.font.size = Pt(10)

# ===== プレイヤーボックス =====
# ダイヤゼブラ電機（左）
dz_box = add_rounded_rect(slide, 1.5, 6, 6, 3, RGBColor(200, 230, 255),
    "ダイヤゼブラ電機\n（旧 田淵電機）", font_size=13, font_color=DARK_BLUE)

dz_sub = add_text_box(slide, 1.5, 9.2, 6, 1.2,
    "最新PCS製造・供給", font_size=9, color=GRAY)

# ユアサ商事（中央）
yuasa_box = add_rounded_rect(slide, 11.5, 6, 6.5, 3, RGBColor(255, 230, 200),
    "ユアサ商事", font_size=14, font_color=RGBColor(180, 90, 0))

yuasa_sub = add_text_box(slide, 11.5, 9.2, 6.5, 1.5,
    "PCS在庫確保\nファイナンス（与信）提供", font_size=9, color=GRAY)

# オルテナジー（右上）
alt_box = add_rounded_rect(slide, 21, 6, 6.5, 3, RGBColor(200, 255, 200),
    "オルテナジー", font_size=14, font_color=GREEN)

alt_sub = add_text_box(slide, 21, 9.2, 6.5, 1.2,
    "「パワまる」実装・販売", font_size=9, color=GRAY)

# ユーザ（右下）
user_box = add_rounded_rect(slide, 21, 13, 6.5, 2.5, RGBColor(255, 255, 200),
    "FIT発電所オーナー\n（エンドユーザ）", font_size=12, font_color=RGBColor(150, 120, 0))

# ===== フロー矢印と説明テキスト =====

# 矢印1: ダイヤゼブラ → ユアサ
arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(7.7), Cm(7), Cm(3.5), Cm(1))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = LIGHT_BLUE
arrow1.line.fill.background()

add_text_box(slide, 7.7, 8.2, 3.5, 1, "①PCS供給", font_size=9, bold=True, color=DARK_BLUE)

# 矢印2: ユアサ → オルテナジー
arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(18.2), Cm(7), Cm(2.5), Cm(1))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = ORANGE
arrow2.line.fill.background()

add_text_box(slide, 18, 8.2, 3, 1, "②在庫提供", font_size=9, bold=True, color=RGBColor(180, 90, 0))

# 矢印3: オルテナジー → ユーザ
arrow3 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Cm(23.5), Cm(9.5), Cm(1), Cm(3.2))
arrow3.fill.solid()
arrow3.fill.fore_color.rgb = GREEN
arrow3.line.fill.background()

add_text_box(slide, 24.8, 10.5, 3.5, 1.5, "③「パワまる」\n提供", font_size=9, bold=True, color=GREEN)

# 矢印4: ユアサ → ユーザ（ファイナンス）- 点線風に
arrow4 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Cm(14), Cm(9.3), Cm(0.8), Cm(5))
arrow4.fill.solid()
arrow4.fill.fore_color.rgb = RGBColor(255, 180, 100)
arrow4.line.color.rgb = ORANGE
arrow4.rotation = 30

add_text_box(slide, 15, 11.5, 4, 1.5, "④ファイナンス\n（与信提供）", font_size=9, bold=True, color=RGBColor(180, 90, 0))

# ===== パワまる説明ボックス =====
pawa_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(13), Cm(12), Cm(4))
pawa_box.fill.solid()
pawa_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
pawa_box.line.color.rgb = GREEN

pawa_text = pawa_box.text_frame
pawa_text.word_wrap = True

pp1 = pawa_text.paragraphs[0]
pp1.text = "【パワまる サービス概要】"
pp1.font.size = Pt(11)
pp1.font.bold = True
pp1.font.color.rgb = GREEN

pp2 = pawa_text.add_paragraph()
pp2.text = "• 最新型PCSへの交換＋付帯サービスのセット"
pp2.font.size = Pt(9)

pp3 = pawa_text.add_paragraph()
pp3.text = "• 月額16,800円～（初期費用・頭金不要）"
pp3.font.size = Pt(9)

pp4 = pawa_text.add_paragraph()
pp4.text = "• 監視代行・駆けつけ対応・故障時無償交換"
pp4.font.size = Pt(9)

pp5 = pawa_text.add_paragraph()
pp5.text = "• 契約満了後はPCS無償譲渡"
pp5.font.size = Pt(9)

pp6 = pawa_text.add_paragraph()
pp6.text = "※本件ではダイヤゼブラ製PCSを採用"
pp6.font.size = Pt(9)
pp6.font.bold = True
pp6.font.color.rgb = DARK_BLUE

# ===== ポイント説明 =====
point_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(17.3), Cm(27.7), Cm(3.2))
point_box.fill.solid()
point_box.fill.fore_color.rgb = RGBColor(255, 250, 230)
point_box.line.color.rgb = ORANGE

pt_text = point_box.text_frame
pt_text.word_wrap = True

pt1 = pt_text.paragraphs[0]
pt1.text = "【本スキームのポイント】"
pt1.font.size = Pt(11)
pt1.font.bold = True
pt1.font.color.rgb = RGBColor(180, 90, 0)

pt2 = pt_text.add_paragraph()
pt2.text = "① ユアサ商事がダイヤゼブラから最新PCSを仕入れ、在庫として確保（安定供給体制の構築）"
pt2.font.size = Pt(9)

pt3 = pt_text.add_paragraph()
pt3.text = "② オルテナジーがユーザに対し「パワまる」サービスとしてPCS交換ソリューションを提供（施工・保守・監視を一括対応）"
pt3.font.size = Pt(9)

pt4 = pt_text.add_paragraph()
pt4.text = "③ ユアサ商事がエンドユーザの与信を手当てし、ファイナンス面をサポート（月額払いの実現）"
pt4.font.size = Pt(9)

# 保存
output_path = r'C:\Users\田中　圭亮\Desktop\Claude_Code_Demo\商流図_FIT発電所PCS交換ソリューション.pptx'
prs.save(output_path)
print(f"PowerPointファイルを作成しました: {output_path}")
